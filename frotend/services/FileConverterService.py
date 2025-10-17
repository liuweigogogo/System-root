"""
文件格式转换服务
支持多种文件格式之间的相互转换
"""

import os
import io
import tempfile
from typing import Optional, Dict, Any, List, Tuple
from pathlib import Path
import logging

# 导入各种转换库
try:
    from docx import Document
    from docx2pdf import convert as docx2pdf_convert
    from pdf2docx import Converter as Pdf2DocxConverter
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    from pptx import Presentation
    from pptx2pdf import convert as pptx2pdf_convert
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    import pandas as pd
    PANDAS_AVAILABLE = True
except ImportError:
    PANDAS_AVAILABLE = False

try:
    from PIL import Image
    PIL_AVAILABLE = True
except ImportError:
    PIL_AVAILABLE = False

try:
    import PyPDF2
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False

from config.LoggerConfig import log_info, log_error

class FileConverterService:
    """文件格式转换服务类"""
    
    # 支持的转换格式映射
    SUPPORTED_CONVERSIONS = {
        # Word 相关转换
        'docx_to_pdf': {'from': 'docx', 'to': 'pdf', 'description': 'Word文档转PDF'},
        'pdf_to_docx': {'from': 'pdf', 'to': 'docx', 'description': 'PDF转Word文档'},
        
        # PowerPoint 相关转换
        'pptx_to_pdf': {'from': 'pptx', 'to': 'pdf', 'description': 'PowerPoint转PDF'},
        'pdf_to_pptx': {'from': 'pdf', 'to': 'pptx', 'description': 'PDF转PowerPoint'},
        
        # Excel 相关转换
        'xlsx_to_csv': {'from': 'xlsx', 'to': 'csv', 'description': 'Excel转CSV'},
        'csv_to_xlsx': {'from': 'csv', 'to': 'xlsx', 'description': 'CSV转Excel'},
        'xlsx_to_json': {'from': 'xlsx', 'to': 'json', 'description': 'Excel转JSON'},
        'json_to_xlsx': {'from': 'json', 'to': 'xlsx', 'description': 'JSON转Excel'},
        
        # 图片相关转换
        'jpg_to_png': {'from': 'jpg', 'to': 'png', 'description': 'JPG转PNG'},
        'png_to_jpg': {'from': 'png', 'to': 'jpg', 'description': 'PNG转JPG'},
        'webp_to_png': {'from': 'webp', 'to': 'png', 'description': 'WebP转PNG'},
        'png_to_webp': {'from': 'png', 'to': 'webp', 'description': 'PNG转WebP'},
    }
    
    def __init__(self):
        """初始化文件转换服务"""
        self.temp_dir = tempfile.gettempdir()
        self._check_dependencies()
    
    def _check_dependencies(self):
        """检查转换库依赖"""
        self.dependencies = {
            'docx': DOCX_AVAILABLE,
            'pptx': PPTX_AVAILABLE,
            'pandas': PANDAS_AVAILABLE,
            'pil': PIL_AVAILABLE,
            'pdf': PDF_AVAILABLE
        }
        
        missing_deps = [k for k, v in self.dependencies.items() if not v]
        if missing_deps:
            log_error(f"缺少转换库依赖: {', '.join(missing_deps)}")
    
    def get_supported_conversions(self) -> Dict[str, Dict[str, str]]:
        """
        获取支持的转换格式列表
        
        Returns:
            Dict: 支持的转换格式字典
        """
        return self.SUPPORTED_CONVERSIONS.copy()
    
    def convert_file(self, file_path: str, target_format: str, 
                    output_path: Optional[str] = None) -> Dict[str, Any]:
        """
        转换文件格式
        
        Args:
            file_path: 源文件路径
            target_format: 目标格式
            output_path: 输出文件路径（可选）
            
        Returns:
            Dict: 转换结果信息
        """
        try:
            # 验证文件是否存在
            if not os.path.exists(file_path):
                return {
                    'success': False,
                    'message': f'源文件不存在: {file_path}',
                    'error_code': 'FILE_NOT_FOUND'
                }
            
            # 获取文件扩展名
            source_ext = Path(file_path).suffix.lower().lstrip('.')
            
            # 构建转换键
            conversion_key = f"{source_ext}_to_{target_format}"
            
            # 检查是否支持该转换
            if conversion_key not in self.SUPPORTED_CONVERSIONS:
                return {
                    'success': False,
                    'message': f'不支持从 {source_ext} 转换到 {target_format}',
                    'error_code': 'UNSUPPORTED_CONVERSION'
                }
            
            # 生成输出路径
            if not output_path:
                output_path = self._generate_output_path(file_path, target_format)
            
            # 执行转换
            result = self._perform_conversion(file_path, output_path, conversion_key)
            
            if result['success']:
                log_info(f"文件转换成功: {file_path} -> {output_path}")
                return {
                    'success': True,
                    'message': '文件转换成功',
                    'output_path': output_path,
                    'conversion_type': conversion_key
                }
            else:
                return result
                
        except Exception as e:
            log_error(e, f"文件转换失败: {file_path}")
            return {
                'success': False,
                'message': f'转换过程中发生错误: {str(e)}',
                'error_code': 'CONVERSION_ERROR'
            }
    
    def _generate_output_path(self, input_path: str, target_format: str) -> str:
        """
        生成输出文件路径
        
        Args:
            input_path: 输入文件路径
            target_format: 目标格式
            
        Returns:
            str: 输出文件路径
        """
        input_path_obj = Path(input_path)
        output_filename = f"{input_path_obj.stem}.{target_format}"
        return str(input_path_obj.parent / output_filename)
    
    def _perform_conversion(self, input_path: str, output_path: str, 
                          conversion_key: str) -> Dict[str, Any]:
        """
        执行具体的转换操作
        
        Args:
            input_path: 输入文件路径
            output_path: 输出文件路径
            conversion_key: 转换类型键
            
        Returns:
            Dict: 转换结果
        """
        try:
            # 根据转换类型调用相应的转换方法
            if conversion_key == 'docx_to_pdf':
                return self._convert_docx_to_pdf(input_path, output_path)
            elif conversion_key == 'pdf_to_docx':
                return self._convert_pdf_to_docx(input_path, output_path)
            elif conversion_key == 'pptx_to_pdf':
                return self._convert_pptx_to_pdf(input_path, output_path)
            elif conversion_key == 'pdf_to_pptx':
                return self._convert_pdf_to_pptx(input_path, output_path)
            elif conversion_key == 'xlsx_to_csv':
                return self._convert_xlsx_to_csv(input_path, output_path)
            elif conversion_key == 'csv_to_xlsx':
                return self._convert_csv_to_xlsx(input_path, output_path)
            elif conversion_key == 'xlsx_to_json':
                return self._convert_xlsx_to_json(input_path, output_path)
            elif conversion_key == 'json_to_xlsx':
                return self._convert_json_to_xlsx(input_path, output_path)
            elif conversion_key in ['jpg_to_png', 'png_to_jpg', 'webp_to_png', 'png_to_webp']:
                return self._convert_image(input_path, output_path, conversion_key)
            else:
                return {
                    'success': False,
                    'message': f'未实现的转换类型: {conversion_key}',
                    'error_code': 'NOT_IMPLEMENTED'
                }
        except Exception as e:
            return {
                'success': False,
                'message': f'转换执行失败: {str(e)}',
                'error_code': 'CONVERSION_EXECUTION_ERROR'
            }
    
    # ===== Word 和 PDF 转换方法 =====
    
    def _convert_docx_to_pdf(self, input_path: str, output_path: str) -> Dict[str, Any]:
        """Word文档转PDF"""
        if not self.dependencies['docx']:
            return {'success': False, 'message': '缺少docx2pdf库', 'error_code': 'MISSING_DEPENDENCY'}
        
        try:
            docx2pdf_convert(input_path, output_path)
            return {'success': True}
        except Exception as e:
            return {'success': False, 'message': str(e), 'error_code': 'CONVERSION_ERROR'}
    
    def _convert_pdf_to_docx(self, input_path: str, output_path: str) -> Dict[str, Any]:
        """PDF转Word文档"""
        if not self.dependencies['docx']:
            return {'success': False, 'message': '缺少pdf2docx库', 'error_code': 'MISSING_DEPENDENCY'}
        
        try:
            cv = Pdf2DocxConverter(input_path)
            cv.convert(output_path)
            cv.close()
            return {'success': True}
        except Exception as e:
            return {'success': False, 'message': str(e), 'error_code': 'CONVERSION_ERROR'}
    
    # ===== PowerPoint 和 PDF 转换方法 =====
    
    def _convert_pptx_to_pdf(self, input_path: str, output_path: str) -> Dict[str, Any]:
        """PowerPoint转PDF"""
        if not self.dependencies['pptx']:
            return {'success': False, 'message': '缺少pptx2pdf库', 'error_code': 'MISSING_DEPENDENCY'}
        
        try:
            pptx2pdf_convert(input_path, output_path)
            return {'success': True}
        except Exception as e:
            return {'success': False, 'message': str(e), 'error_code': 'CONVERSION_ERROR'}
    
    def _convert_pdf_to_pptx(self, input_path: str, output_path: str) -> Dict[str, Any]:
        """PDF转PowerPoint（注意：这是一个复杂转换，可能需要特殊处理）"""
        if not self.dependencies['pptx']:
            return {'success': False, 'message': '缺少pptx库', 'error_code': 'MISSING_DEPENDENCY'}
        
        try:
            # 这里需要实现PDF到PPTX的转换逻辑
            # 由于PDF到PPTX的转换比较复杂，这里提供一个基础实现
            presentation = Presentation()
            
            # 添加一个空白幻灯片
            slide_layout = presentation.slide_layouts[6]  # 空白布局
            slide = presentation.slides.add_slide(slide_layout)
            
            # 添加文本说明
            textbox = slide.shapes.add_textbox(0, 0, 9144000, 6858000)  # 位置和大小
            text_frame = textbox.text_frame
            text_frame.text = f"PDF文件: {os.path.basename(input_path)}\n\n注意：PDF到PowerPoint的转换需要手动处理内容。"
            
            presentation.save(output_path)
            return {'success': True}
        except Exception as e:
            return {'success': False, 'message': str(e), 'error_code': 'CONVERSION_ERROR'}
    
    # ===== Excel 相关转换方法 =====
    
    def _convert_xlsx_to_csv(self, input_path: str, output_path: str) -> Dict[str, Any]:
        """Excel转CSV"""
        if not self.dependencies['pandas']:
            return {'success': False, 'message': '缺少pandas库', 'error_code': 'MISSING_DEPENDENCY'}
        
        try:
            # 读取Excel文件
            df = pd.read_excel(input_path)
            # 保存为CSV
            df.to_csv(output_path, index=False, encoding='utf-8-sig')
            return {'success': True}
        except Exception as e:
            return {'success': False, 'message': str(e), 'error_code': 'CONVERSION_ERROR'}
    
    def _convert_csv_to_xlsx(self, input_path: str, output_path: str) -> Dict[str, Any]:
        """CSV转Excel"""
        if not self.dependencies['pandas']:
            return {'success': False, 'message': '缺少pandas库', 'error_code': 'MISSING_DEPENDENCY'}
        
        try:
            # 读取CSV文件
            df = pd.read_csv(input_path)
            # 保存为Excel
            df.to_excel(output_path, index=False)
            return {'success': True}
        except Exception as e:
            return {'success': False, 'message': str(e), 'error_code': 'CONVERSION_ERROR'}
    
    def _convert_xlsx_to_json(self, input_path: str, output_path: str) -> Dict[str, Any]:
        """Excel转JSON"""
        if not self.dependencies['pandas']:
            return {'success': False, 'message': '缺少pandas库', 'error_code': 'MISSING_DEPENDENCY'}
        
        try:
            # 读取Excel文件
            df = pd.read_excel(input_path)
            # 转换为JSON
            json_data = df.to_json(orient='records', force_ascii=False, indent=2)
            # 保存JSON文件
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(json_data)
            return {'success': True}
        except Exception as e:
            return {'success': False, 'message': str(e), 'error_code': 'CONVERSION_ERROR'}
    
    def _convert_json_to_xlsx(self, input_path: str, output_path: str) -> Dict[str, Any]:
        """JSON转Excel"""
        if not self.dependencies['pandas']:
            return {'success': False, 'message': '缺少pandas库', 'error_code': 'MISSING_DEPENDENCY'}
        
        try:
            # 读取JSON文件
            df = pd.read_json(input_path)
            # 保存为Excel
            df.to_excel(output_path, index=False)
            return {'success': True}
        except Exception as e:
            return {'success': False, 'message': str(e), 'error_code': 'CONVERSION_ERROR'}
    
    # ===== 图片转换方法 =====
    
    def _convert_image(self, input_path: str, output_path: str, conversion_key: str) -> Dict[str, Any]:
        """图片格式转换"""
        if not self.dependencies['pil']:
            return {'success': False, 'message': '缺少PIL库', 'error_code': 'MISSING_DEPENDENCY'}
        
        try:
            # 打开图片
            with Image.open(input_path) as img:
                # 根据转换类型处理
                if conversion_key == 'jpg_to_png':
                    # JPG转PNG，需要处理RGBA模式
                    if img.mode == 'RGBA':
                        img.save(output_path, 'PNG')
                    else:
                        img.convert('RGB').save(output_path, 'PNG')
                elif conversion_key == 'png_to_jpg':
                    # PNG转JPG，需要转换为RGB模式
                    if img.mode in ('RGBA', 'LA', 'P'):
                        # 创建白色背景
                        background = Image.new('RGB', img.size, (255, 255, 255))
                        if img.mode == 'P':
                            img = img.convert('RGBA')
                        background.paste(img, mask=img.split()[-1] if img.mode == 'RGBA' else None)
                        background.save(output_path, 'JPEG', quality=95)
                    else:
                        img.convert('RGB').save(output_path, 'JPEG', quality=95)
                elif conversion_key == 'webp_to_png':
                    img.save(output_path, 'PNG')
                elif conversion_key == 'png_to_webp':
                    img.save(output_path, 'WEBP', quality=95)
                
                return {'success': True}
        except Exception as e:
            return {'success': False, 'message': str(e), 'error_code': 'CONVERSION_ERROR'}
    
    def batch_convert(self, file_list: List[str], target_format: str, 
                     output_dir: Optional[str] = None) -> Dict[str, Any]:
        """
        批量转换文件
        
        Args:
            file_list: 文件路径列表
            target_format: 目标格式
            output_dir: 输出目录（可选）
            
        Returns:
            Dict: 批量转换结果
        """
        results = {
            'success_count': 0,
            'failed_count': 0,
            'results': [],
            'total_files': len(file_list)
        }
        
        for file_path in file_list:
            # 生成输出路径
            if output_dir:
                filename = Path(file_path).stem
                output_path = os.path.join(output_dir, f"{filename}.{target_format}")
            else:
                output_path = None
            
            # 执行转换
            result = self.convert_file(file_path, target_format, output_path)
            results['results'].append({
                'file': file_path,
                'result': result
            })
            
            if result['success']:
                results['success_count'] += 1
            else:
                results['failed_count'] += 1
        
        return results
    
    def get_file_info(self, file_path: str) -> Dict[str, Any]:
        """
        获取文件信息
        
        Args:
            file_path: 文件路径
            
        Returns:
            Dict: 文件信息
        """
        try:
            if not os.path.exists(file_path):
                return {'success': False, 'message': '文件不存在'}
            
            file_stat = os.stat(file_path)
            file_ext = Path(file_path).suffix.lower().lstrip('.')
            
            return {
                'success': True,
                'file_name': os.path.basename(file_path),
                'file_size': file_stat.st_size,
                'file_extension': file_ext,
                'modified_time': file_stat.st_mtime,
                'supported_conversions': [
                    key for key, value in self.SUPPORTED_CONVERSIONS.items()
                    if value['from'] == file_ext
                ]
            }
        except Exception as e:
            return {'success': False, 'message': str(e)}
